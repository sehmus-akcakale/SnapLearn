"""
PowerPoint Slide Generator Module
=================================
This module creates educational slides using the python-pptx library.
A new presentation file is created for each session, and all screenshots
during that session are added to the same file.
"""

import logging
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from config import (
    OUTPUT_DIR,
    SLIDE_WIDTH_INCHES,
    SLIDE_HEIGHT_INCHES,
    IMAGE_WIDTH_INCHES,
    IMAGE_HEIGHT_INCHES,
    TEXT_BOX_WIDTH_INCHES,
    TEXT_BOX_HEIGHT_INCHES
)

# Setup logger
logger = logging.getLogger(__name__)


class SlideGenerator:
    """
    PowerPoint presentation generator class.
    
    Creates a unique presentation file for each session and
    manages adding new slides to it.
    """
    
    def __init__(self):
        """
        Creates a new SlideGenerator instance.
        Automatically starts a new presentation.
        """
        self.presentation = None
        self.filepath = None
        self.slide_count = 0
        self._create_new_presentation()
    
    def _create_new_presentation(self):
        """
        Creates a new empty presentation with a session-specific filename.
        """
        # Create new presentation
        self.presentation = Presentation()
        
        # Set slide dimensions (16:9 widescreen)
        self.presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
        self.presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)
        
        # Create unique filename for this session
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        self.filepath = OUTPUT_DIR / f"presentation_{timestamp}.pptx"
        
        # Add title slide
        self._add_title_slide()
        
        logger.info(f"New presentation created: {self.filepath}")
    
    def _add_title_slide(self):
        """
        Adds a title slide to the presentation.
        """
        # Use blank slide layout (index 6 is typically blank)
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title text box
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(SLIDE_WIDTH_INCHES - 1)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Educational Notes"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)  # Blue
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_top = Inches(4.0)
        subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        
        date_str = datetime.now().strftime("%B %d, %Y")
        subtitle_para.text = f"Auto-generated - {date_str}"
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)  # Gray
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        logger.info("Title slide added.")
    
    def add_content_slide(
        self, 
        image_path: Path, 
        summary: str, 
        question: str,
        slide_title: str = None
    ) -> int:
        """
        Adds two slides: first with full-screen image, second with text content.
        
        Args:
            image_path: Path to the image file to add
            summary: AI-generated summary text
            question: AI-generated multiple choice question
            slide_title: Slide title (optional)
            
        Returns:
            int: The slide number that was added
        """
        try:
            self.slide_count += 1
            
            # Slide title
            if slide_title is None:
                slide_title = f"Slide {self.slide_count}"
            
            # === SLIDE 1: Full-screen image ===
            slide_layout = self.presentation.slide_layouts[6]
            image_slide = self.presentation.slides.add_slide(slide_layout)
            
            self._add_slide_title(image_slide, slide_title)
            self._add_full_image(image_slide, image_path)
            
            # === SLIDE 2: Summary and Question ===
            text_slide = self.presentation.slides.add_slide(slide_layout)
            
            self._add_slide_title(text_slide, f"{slide_title} - Notes")
            self._add_full_text_content(text_slide, summary, question)
            
            # Save presentation
            self._save()
            
            logger.info(f"Content slides #{self.slide_count} added (image + text).")
            
            return self.slide_count
            
        except Exception as e:
            logger.error(f"Error adding slide: {e}")
            raise
    
    def _add_slide_title(self, slide, title: str):
        """
        Adds a title to the slide.
        """
        left = Inches(0.3)
        top = Inches(0.2)
        width = Inches(SLIDE_WIDTH_INCHES - 0.6)
        height = Inches(0.6)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)  # Dark blue
    
    def _add_image(self, slide, image_path: Path):
        """
        Adds an image to the left side of the slide.
        """
        # Image position (left side)
        left = Inches(0.3)
        top = Inches(1.0)
        
        # Add and resize image
        picture = slide.shapes.add_picture(
            str(image_path),
            left,
            top,
            width=Inches(IMAGE_WIDTH_INCHES),
            height=Inches(IMAGE_HEIGHT_INCHES)
        )
        
        logger.debug(f"Image added: {image_path}")
    
    def _add_full_text_content(self, slide, summary: str, question: str):
        """
        Adds summary and question text as full-width content on a dedicated slide.
        """
        # Text box position (full width, centered)
        left = Inches(0.5)
        top = Inches(1.0)
        width = Inches(SLIDE_WIDTH_INCHES - 1.0)
        height = Inches(6.0)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Summary header
        p1 = text_frame.paragraphs[0]
        p1.text = "Summary"
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
        p1.space_after = Pt(12)
        
        # Summary text
        p2 = text_frame.add_paragraph()
        # Allow more characters since we have full width
        max_chars = 800
        display_summary = summary[:max_chars] + "..." if len(summary) > max_chars else summary
        p2.text = display_summary
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p2.space_after = Pt(24)
        
        # Question header
        p3 = text_frame.add_paragraph()
        p3.text = "Multiple Choice Question"
        p3.font.size = Pt(24)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(0xC0, 0x50, 0x4D)  # Red
        p3.space_after = Pt(12)
        
        # Question text
        p4 = text_frame.add_paragraph()
        p4.text = question
        p4.font.size = Pt(16)
        p4.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
        logger.debug("Full text content added.")
    
    def add_direct_slide(self, image_path: Path, slide_title: str = None) -> int:
        """
        Adds a slide with only the screenshot (full width, no text).
        Used for direct capture without AI analysis.
        
        Args:
            image_path: Path to the image file to add
            slide_title: Slide title (optional)
            
        Returns:
            int: The slide number that was added
        """
        try:
            self.slide_count += 1
            
            # Use blank slide layout
            slide_layout = self.presentation.slide_layouts[6]
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Slide title
            if slide_title is None:
                slide_title = f"Slide {self.slide_count} (Direct Capture)"
            
            self._add_slide_title(slide, slide_title)
            
            # Add image (larger, centered)
            self._add_full_image(slide, image_path)
            
            # Save presentation
            self._save()
            
            logger.info(f"Direct capture slide #{self.slide_count} added.")
            
            return self.slide_count
            
        except Exception as e:
            logger.error(f"Error adding direct slide: {e}")
            raise
    
    def _add_full_image(self, slide, image_path: Path):
        """
        Adds a larger image centered on the slide (for direct capture).
        """
        # Image position (centered, larger)
        img_width = Inches(11.0)
        img_height = Inches(5.8)
        left = Inches((SLIDE_WIDTH_INCHES - 11.0) / 2)
        top = Inches(1.0)
        
        # Add and resize image
        picture = slide.shapes.add_picture(
            str(image_path),
            left,
            top,
            width=img_width,
            height=img_height
        )
        
        logger.debug(f"Full image added: {image_path}")
    
    def _save(self):
        """
        Saves the presentation to file.
        """
        try:
            self.presentation.save(str(self.filepath))
            logger.debug(f"Presentation saved: {self.filepath}")
        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
            raise
    
    def get_filepath(self) -> Path:
        """
        Returns the current presentation file path.
        """
        return self.filepath
    
    def get_slide_count(self) -> int:
        """
        Returns the number of content slides added (excluding title slide).
        """
        return self.slide_count


# =============================================================================
# Simple usage functions
# =============================================================================

# Global generator instance (used throughout the session)
_generator: SlideGenerator = None


def get_generator() -> SlideGenerator:
    """
    Returns the SlideGenerator instance for the current session.
    Creates a new instance on first call.
    """
    global _generator
    if _generator is None:
        _generator = SlideGenerator()
    return _generator


def add_slide(image_path: Path, summary: str, question: str) -> int:
    """
    Adds a new slide to the current presentation.
    
    Args:
        image_path: Path to the image file to add
        summary: Summary text
        question: Multiple choice question
        
    Returns:
        int: The slide number that was added
    """
    generator = get_generator()
    return generator.add_content_slide(image_path, summary, question)


def add_slide_direct(image_path: Path) -> int:
    """
    Adds a new slide with only the screenshot (no AI analysis).
    
    Args:
        image_path: Path to the image file to add
        
    Returns:
        int: The slide number that was added
    """
    generator = get_generator()
    return generator.add_direct_slide(image_path)


def get_current_filepath() -> Path:
    """
    Returns the current presentation file path.
    """
    generator = get_generator()
    return generator.get_filepath()


def reset_session():
    """
    Starts a new session (creates a new presentation file).
    """
    global _generator
    _generator = SlideGenerator()
    return _generator.get_filepath()


def export_to_pdf() -> Path | None:
    """
    Exports the current presentation to PDF format.
    Uses PowerPoint COM automation on Windows.
    
    Returns:
        Path | None: Path to the PDF file or None on error
    """
    generator = get_generator()
    pptx_path = generator.get_filepath()
    
    if not pptx_path.exists():
        logger.error("Presentation file does not exist!")
        return None
    
    pdf_path = pptx_path.with_suffix('.pdf')
    
    try:
        import subprocess
        import sys
        
        if sys.platform == "win32":
            # Use PowerPoint COM automation via PowerShell
            ps_script = f'''
$pptx = "{pptx_path.absolute()}"
$pdf = "{pdf_path.absolute()}"
$powerpoint = New-Object -ComObject PowerPoint.Application
$powerpoint.Visible = [Microsoft.Office.Core.MsoTriState]::msoFalse
$presentation = $powerpoint.Presentations.Open($pptx, [Microsoft.Office.Core.MsoTriState]::msoTrue, [Microsoft.Office.Core.MsoTriState]::msoFalse, [Microsoft.Office.Core.MsoTriState]::msoFalse)
$presentation.SaveAs($pdf, 32)
$presentation.Close()
$powerpoint.Quit()
'''
            # Run PowerShell script
            result = subprocess.run(
                ["powershell", "-Command", ps_script],
                capture_output=True,
                text=True,
                timeout=60
            )
            
            if result.returncode == 0 and pdf_path.exists():
                logger.info(f"PDF exported: {pdf_path}")
                return pdf_path
            else:
                logger.error(f"PDF export failed: {result.stderr}")
                return None
        else:
            logger.warning("PDF export is only supported on Windows with PowerPoint installed.")
            return None
            
    except subprocess.TimeoutExpired:
        logger.error("PDF export timed out!")
        return None
    except Exception as e:
        logger.error(f"PDF export error: {e}")
        return None


# =============================================================================
# Test / Demo
# =============================================================================

if __name__ == "__main__":
    # Setup logging
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s - %(name)s - %(levelname)s - %(message)s"
    )
    
    print("Slide generator test...")
    
    # Create test presentation
    generator = SlideGenerator()
    
    print(f"Presentation file: {generator.get_filepath()}")
    print("Test successful! (A real image file is needed to add content slides)")
